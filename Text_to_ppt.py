import streamlit as st
import pptx
from pptx.util import Pt
import wikipediaapi
import os
import base64

# Custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(18)
MAX_SLIDES = 12  # Maximum number of slides

# Function to fetch Wikipedia content
def fetch_wikipedia_content(topic):
    wiki_wiki = wikipediaapi.Wikipedia(
        language='en',
        user_agent="TextToPPTApp/1.0 (contact: smitgevariya.sg141@gmail.com) WikipediaAPI/0.5.8"
    )
    page = wiki_wiki.page(topic)
    if page.exists():
        content = page.text.split('\n\n')  # Split by paragraphs
        content = [paragraph.strip() for paragraph in content if paragraph.strip()]
        return content[:MAX_SLIDES]  # Limit to first 12 sections
    else:
        return ["No information found for the given topic. Please try another topic."]

# Function to condense content into bullet points
def condense_content(content, max_slides):
    condensed_content = []
    for paragraph in content:
        points = paragraph.split('. ')  # Split by sentences
        condensed_content.extend(points[:3])  # Take first 3 sentences as bullet points
        if len(condensed_content) >= max_slides * 3:  # Ensure slide count limit
            break
    return [condensed_content[i:i + 3] for i in range(0, len(condensed_content), 3)][:max_slides]

# Function to generate slide titles and content
def generate_slide_titles_and_content(content):
    slide_titles = []
    for paragraph in content:
        # Use the first line of each paragraph as the slide title
        title = paragraph.split('.')[0]  # Get the first sentence
        slide_titles.append(title)
    return slide_titles, condense_content(content, MAX_SLIDES)

# Function to create a PowerPoint presentation
def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    # Add content slides
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title  # Set title based on content
        slide.placeholders[1].text = '\n'.join(slide_content)  # Join bullet points

        # Apply custom font sizes
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    # Add Thank You slide
    thank_you_slide = prs.slides.add_slide(prs.slide_layouts[0])
    thank_you_slide.shapes.title.text = "Thank You"
    thank_you_slide.placeholders[1].text = "Questions?"

    # Save presentation
    output_dir = "generated_ppt"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"{topic}_presentation.pptx")
    prs.save(output_path)
    return output_path

# Function to create a download link for the PPT
def get_ppt_download_link(ppt_path):
    with open(ppt_path, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{os.path.basename(ppt_path)}">Download your presentation</a>'

# Main function for Streamlit app
def main():
    st.title("Text to PowerPoint Generator")
    st.write("Generate a PowerPoint presentation from a Wikipedia topic.")

    # Input topic
    topic = st.text_input("Enter a topic to search on Wikipedia:")

    if st.button("Generate PPT"):
        if topic.strip():
            st.write("Fetching content and generating presentation...")
            content = fetch_wikipedia_content(topic)
            slide_titles, slide_contents = generate_slide_titles_and_content(content)
            ppt_path = create_presentation(topic, slide_titles, slide_contents)
            st.success("Presentation generated successfully!")
            st.markdown(get_ppt_download_link(ppt_path), unsafe_allow_html=True)
        else:
            st.error("Please enter a valid topic.")

if __name__ == "__main__":
    main()
